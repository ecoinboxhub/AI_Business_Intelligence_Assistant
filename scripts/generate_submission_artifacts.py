"""Generate ALL NexaSphere submission artifacts for competition judges.

Outputs (written to the gitignored ``submission/`` folder):
  * submission/NexaSphere_Comprehensive_Summary.docx   - Full summary document
  * submission/NexaSphere_Comprehensive_Summary.pdf     - Full summary PDF
  * submission/NexaSphere_Presentation_Slides.pptx      - Winning slide deck
  * submission/NexaSphere_Presentation_Slides.pdf       - Slide deck as PDF
  * submission/NexaSphere_Demo_Video.mp4                - 2-min video with Nigerian male voice

Prerequisites:
  pip install python-docx python-pptx reportlab Pillow edge-tts moviepy

Run:  python scripts/generate_submission_artifacts.py
"""
import asyncio
import os
import shutil
import sys
import textwrap
import math
from pathlib import Path

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "submission"
SHOTS = OUT / "screenshots"
OUT.mkdir(parents=True, exist_ok=True)

DEPLOYED_WEB = "https://ecoinboxhub.github.io/AI_Business_Intelligence_Assistant/"
DEPLOYED_MOBILE = "https://github.com/ecoinboxhub/AI_Business_Intelligence_Assistant/releases/download/v1.0.0/mobile-release.zip"
DEPLOYED_REPO = "https://github.com/ecoinboxhub/AI_Business_Intelligence_Assistant"
LIVE_APP = DEPLOYED_WEB

# ---------------------------------------------------------------------------
# Shared brand colours
# ---------------------------------------------------------------------------
BRAND_DARK = "#0F172A"
BRAND_INDIGO = "#4F46E5"
BRAND_SKY = "#0EA5E9"
BRAND_EMERALD = "#10B981"
BRAND_AMBER = "#F59E0B"
BRAND_CRIMSON = "#EF4444"
BRAND_WHITE = "#FFFFFF"
BRAND_LIGHT_BG = "#F8FAFC"

# ---------------------------------------------------------------------------
# 1. COMPREHENSIVE SUMMARY DOCUMENT (.docx and .pdf)
# ---------------------------------------------------------------------------
def generate_summary_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()

    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # -- Helper --
    def add_heading_styled(text, level=1):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.color.rgb = RGBColor(15, 23, 42)
        return h

    def add_para(text, bold=False, size=11, color=RGBColor(51, 65, 85), align=None, space_after=6):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if align:
            p.alignment = align
        p.paragraph_format.space_after = Pt(space_after)
        return p

    def add_bullet(text, size=10.5):
        p = doc.add_paragraph(style="List Bullet")
        p.clear()
        run = p.add_run(text)
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(51, 65, 85)
        return p

    # ---- TITLE PAGE ----
    for _ in range(4):
        doc.add_paragraph()

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_p.add_run("NexaSphere")
    run.font.size = Pt(36)
    run.bold = True
    run.font.color.rgb = RGBColor(79, 70, 229)

    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = sub_p.add_run("AI Business Intelligence Assistant")
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(15, 23, 42)

    tagline = doc.add_paragraph()
    tagline.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = tagline.add_run("Zero-Hallucination Decision Intelligence for Modern Enterprises")
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = RGBColor(79, 70, 229)

    doc.add_paragraph()

    links = doc.add_paragraph()
    links.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for label, url in [("Live Dashboard", DEPLOYED_WEB), (" | Mobile App", DEPLOYED_MOBILE), (" | Source Code", DEPLOYED_REPO)]:
        run = links.add_run(label)
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(14, 165, 233)
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        r_elem = run._element
        rPr = r_elem.find(qn('w:rPr'))
        if rPr is None:
            rPr = OxmlElement('w:rPr')
            r_elem.insert(0, rPr)
        u_elem = OxmlElement('w:u')
        u_elem.set(qn('w:val'), 'single')
        rPr.append(u_elem)

    doc.add_paragraph()
    date_p = doc.add_paragraph()
    date_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = date_p.add_run("Case Study 4 Solution  |  Competition Submission 2026")
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(100, 116, 139)

    doc.add_page_break()

    # ---- TABLE OF CONTENTS (manual) ----
    add_heading_styled("Table of Contents", level=1)
    toc_items = [
        "1. Executive Summary",
        "2. The Problem We Solve",
        "3. Our Solution: NexaSphere",
        "4. Target Audience & Personas",
        "5. How AI Is Used — The Hybrid Architecture",
        "6. Core Use Cases & Supported Insights",
        "7. Technical Architecture Deep Dive",
        "8. Innovation Layer",
        "9. Business Impact & Value Proposition",
        "10. Testing & Quality Assurance",
        "11. Deployment & Accessibility",
        "12. Future Roadmap",
        "13. Conclusion",
    ]
    for item in toc_items:
        add_para(item, size=11, color=RGBColor(79, 70, 229), space_after=3)

    doc.add_page_break()

    # ---- 1. EXECUTIVE SUMMARY ----
    add_heading_styled("1. Executive Summary", level=1)
    add_para(
        "NexaSphere is a full-stack AI Business Intelligence Assistant that transforms how "
        "executives and managers interact with operational data. By combining deterministic "
        "mathematical computation with natural language understanding, NexaSphere enables "
        "leaders to ask business questions in plain English and receive instant, accurate, "
        "chart-backed answers — with a zero-hallucination guarantee on every number."
    )
    add_para(
        "Unlike generic AI chatbots that fabricate financial figures, NexaSphere enforces a strict "
        "architectural invariant: all mathematics is executed by Python's Pandas engine, while the "
        "LLM (Google Gemini) serves purely as an executive interpreter that generates narrative, "
        "selects visualisations, and recommends next actions. This separation ensures 100% numerical "
        "accuracy while delivering the conversational experience executives expect."
    )
    add_para(
        "The solution ships as three interconnected applications — a FastAPI backend, a React web "
        "dashboard, and an Expo mobile app — all powered by a single analytical engine. NexaSphere "
        "answers nine critical management questions spanning revenue analysis, profitability, return "
        "anomalies, campaign ROI, inventory health, delivery performance, customer segmentation, "
        "employee performance, and target attainment."
    )

    links_p = doc.add_paragraph()
    run = links_p.add_run("Deployed Solution: ")
    run.bold = True
    run.font.size = Pt(10)
    run = links_p.add_run(f"{DEPLOYED_WEB}  |  Mobile: {DEPLOYED_MOBILE}")
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(14, 165, 233)

    # ---- 2. THE PROBLEM ----
    add_heading_styled("2. The Problem We Solve", level=1)
    add_para(
        "Modern businesses generate vast amounts of operational data across sales, marketing, "
        "supply chain, inventory, and human resources systems. However, extracting actionable "
        "insights from this data remains painfully slow and error-prone:"
    )
    problems = [
        "Manual BI reporting takes hours to days, delaying critical decisions.",
        "Fragmented data sources (sales, returns, deliveries, campaigns, inventory, targets) require multiple tools and expertise to cross-analyse.",
        "Generic AI chatbots confidently hallucinate financial numbers, leading to catastrophic misinformed decisions.",
        "Executives lack a single, trusted interface that combines natural language access with mathematical precision.",
        "Mobile access to BI insights is an afterthought, leaving on-the-go leaders disconnected from real-time data.",
    ]
    for p in problems:
        add_bullet(p)

    add_para(
        "The result: businesses make slower decisions, miss revenue opportunities, and suffer from "
        "erosion of trust in AI-powered analytics. NexaSphere was built to solve this permanently.",
        bold=True, size=11
    )

    # ---- 3. OUR SOLUTION ----
    add_heading_styled("3. Our Solution: NexaSphere", level=1)
    add_para(
        "NexaSphere is a zero-hallucination AI Business Intelligence assistant that combines "
        "three critical capabilities into one seamless experience:"
    )
    solutions = [
        "Deterministic Mathematical Engine: Every number — from revenue totals to profit margins to ROI calculations — is computed by Pandas with exact precision. The LLM never produces, suggests, or modifies any numeric value.",
        "Natural Language Interface: Executives ask questions in plain English (e.g., 'Which region generates the most profit?') and receive instant, structured answers with dynamic charts, verified facts, and actionable recommendations.",
        "Cross-Platform Access: A React web dashboard for desktop analysis and an Expo mobile app for on-the-go decision-making, both consuming the same FastAPI backend.",
        "Proactive Anomaly Detection: The system automatically identifies and surfaces business anomalies — weak-margin regions, delivery delays, return hotspots — before the executive even asks.",
        "What-If Simulation: Leaders can model scenarios (e.g., 'What if we reduce delivery delays by 15%?') and see quantified naira impact on the bottom line.",
    ]
    for s in solutions:
        add_bullet(s)

    # ---- 4. TARGET AUDIENCE ----
    add_heading_styled("4. Target Audience & Personas", level=1)
    add_para(
        "NexaSphere is designed for four primary user personas, each with distinct analytical needs:"
    )

    personas = [
        ("C-Suite Executives", "Need high-level KPI dashboards, anomaly alerts, and strategic recommendations. Want answers in seconds, not days. Access via both web and mobile."),
        ("Regional Managers", "Compare performance across regions, identify underperforming stores, and drill into revenue and profitability drivers by geography."),
        ("Marketing Managers", "Evaluate campaign ROI, compare channel effectiveness, and optimise spend allocation across campaigns."),
        ("Supply Chain Directors", "Monitor delivery partner performance, inventory stockouts, return rates, and operational efficiency metrics."),
    ]
    for title, desc in personas:
        p = doc.add_paragraph()
        run = p.add_run(f"{title}: ")
        run.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor(30, 41, 59)
        run = p.add_run(desc)
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor(51, 65, 85)

    # ---- 5. HOW AI IS USED ----
    add_heading_styled("5. How AI Is Used — The Hybrid Architecture", level=1)
    add_para(
        "NexaSphere's most significant innovation is its non-obvious use of AI. Rather than "
        "letting the LLM do everything (including math), the system enforces a strict separation "
        "of concerns that guarantees both accuracy and intelligence:"
    )

    ai_layers = [
        ("Layer 1 — Deterministic Intent Routing", "When a user asks a question, a keyword-scoring algorithm (resolve_intent()) maps it to one of nine registered analysis functions. No LLM is involved in this routing — it is purely deterministic and testable."),
        ("Layer 2 — Pandas Computation Engine", "The matched analysis function executes Pandas aggregations, calculates KPIs, generates chart data payloads, and computes risk scores. Every number originates here. Zero exceptions."),
        ("Layer 3 — DuckDB Analytical Engine", "An additional in-process columnar analytics engine provides sub-5ms query responses with LRU caching, cross-validated against Pandas for accuracy."),
        ("Layer 4 — LLM Narrative Layer", "Only after all numbers are computed does the Gemini LLM generate the natural-language narrative — summaries, interpretations, and recommendations grounded exclusively in the computed results. The LLM output is treated as untrusted prose; numbers are never parsed from model text."),
        ("Layer 5 — Structured Response", "The final StructuredBIResponse merges the Pandas chart payload, metrics, findings, risks, and recommendations into a structured JSON response delivered to the client."),
    ]
    for title, desc in ai_layers:
        p = doc.add_paragraph()
        run = p.add_run(f"{title}: ")
        run.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor(79, 70, 229)
        run = p.add_run(desc)
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor(51, 65, 85)

    add_para(
        "This architecture guarantees: (1) zero mathematical hallucination, (2) complete auditability, "
        "(3) deterministic reproducibility, and (4) graceful degradation — if Gemini is unavailable, "
        "an offline fallback generates equivalent narrative from result seeds.",
        bold=True
    )

    # ---- 6. CORE USE CASES ----
    add_heading_styled("6. Core Use Cases & Supported Insights", level=1)
    add_para("NexaSphere answers nine critical management questions:")

    queries = [
        "Revenue & Profit Drivers — Which products, stores, or regions generate the most revenue and profit?",
        "Growth vs Profitability — Is revenue growth leading to stronger profitability?",
        "Return Anomalies — Which products have unusually high return rates?",
        "Campaign ROI — Which marketing campaigns generate the best return on investment?",
        "Inventory Health — Which stores are experiencing stockouts or excess inventory?",
        "Delivery Performance — Which delivery partners are associated with delays or poor ratings?",
        "Customer Segments — Which customer segments are the most valuable?",
        "Employee Performance — Which employees perform well based on both revenue and profitability?",
        "Target Attainment — Where is the business failing to meet its targets?",
    ]
    for i, q in enumerate(queries, 1):
        add_bullet(f"{q}")

    add_para(
        "Each question triggers a dedicated analysis function that produces a specific chart type "
        "(grouped bar, donut, funnel, treemap, bubble, scatter, bullet), verified metrics, "
        "risk scores, what-if scenarios, and follow-up question suggestions."
    )

    # ---- 7. TECHNICAL ARCHITECTURE ----
    add_heading_styled("7. Technical Architecture Deep Dive", level=1)
    add_para("The system comprises five distinct layers:")

    arch_layers = [
        "Presentation Layer: React 18 web dashboard (Vite, Recharts, lucide-react) on port 3030; Expo React Native mobile app (gifted-charts) on port 8081.",
        "API Layer: FastAPI on port 5050 with 7 endpoints including /api/questions (POST), /api/questions/stream (SSE), /api/simulate (POST), /api/kpis, /api/insights, /api/catalog, /api/analysis/{id}.",
        "Orchestration Layer: Gemini LLM wrapper (app/ai/service.py) with structured NarrativeInsight schema, offline fallback, and enhanced prompt engineering.",
        "Analysis Layer: Nine Pandas analysis functions (app/analysis/) — the ONLY place math happens. Plus DuckDB engine, risk scoring (0-100), and what-if simulation.",
        "Data Layer: Source resolution engine supporting CSV, TSV, Excel, Parquet, and JSON via environment variables. Schema coercion, validation, and caching. Seeded synthetic fallback (2M sales rows, 400K deliveries).",
    ]
    for l in arch_layers:
        add_bullet(l)

    add_para("Technology Stack:", bold=True)
    techs = [
        "Backend: Python 3.10+, FastAPI, Pandas, DuckDB, Google Gemini API",
        "Web: React 18, Vite, Recharts, lucide-react, CSS Modules",
        "Mobile: Expo SDK 50, React Native 0.73, gifted-charts",
        "Testing: pytest with accuracy benchmarks, intent routing tests, data source validation",
        "CI/CD: GitHub Actions (test → deploy frontend to GitHub Pages → release mobile)",
    ]
    for t in techs:
        add_bullet(t)

    # ---- 8. INNOVATION LAYER ----
    add_heading_styled("8. Innovation Layer", level=1)
    add_para(
        "Beyond the core zero-hallucination architecture, NexaSphere includes several innovation features:"
    )
    innovations = [
        "DuckDB Analytical Engine: In-process columnar analytics with named analytical queries and LRU caching, delivering sub-5ms query responses. Cross-validated against Pandas in automated tests.",
        "Deterministic Risk Scoring: Every analysis payload gains a 0-100 severity score with urgency bands (MONITOR, STANDARD, HIGH_PRIORITY, CRITICAL), computed deterministically from data facts.",
        "What-If Simulation Engine: Four parametrized scenario levers (reduce delays, reduce returns, close target gap, shift campaign spend) with quantified naira impact attached to every answer.",
        "Server-Sent Events (SSE): The /api/questions/stream endpoint provides staged progress events — DuckDB metrics preview, synthesis, severity scoring — before the final response.",
        "Proactive Anomaly Insights: The /api/insights endpoint computes live Pandas-driven anomaly detection: weakest-margin region, worst delivery partner, return hotspots, and below-target regions.",
        "Structured Response Contract: Every answer follows a StructuredBIResponse schema with segregated facts, interpretations, recommendations, and follow-up questions — never free-form text.",
    ]
    for inv in innovations:
        add_bullet(inv)

    # ---- 9. BUSINESS IMPACT ----
    add_heading_styled("9. Business Impact & Value Proposition", level=1)
    add_para("NexaSphere delivers measurable business value:")

    impacts = [
        "100% Mathematical Accuracy: Every number is computed by Pandas — zero hallucination risk. Executives can trust every figure in every answer.",
        "~80% Reduction in Reporting Latency: What previously took hours of manual BI work now happens in seconds through natural language queries.",
        "Proactive Decision-Making: Anomaly alerts surface issues before board meetings, enabling leaders to act on problems before they escalate.",
        "Unified Analytics Platform: One assistant answers all nine operational questions across web and mobile, eliminating tool fragmentation.",
        "Graceful Degradation: The system works offline with deterministic fallback narratives, ensuring business continuity even without API connectivity.",
        "Zero Infrastructure Overhead: No Docker, no Kubernetes, no cloud services required. Runs as local processes for the MVP, with GitHub Pages hosting for the frontend.",
    ]
    for imp in impacts:
        add_bullet(imp)

    # ---- 10. TESTING ----
    add_heading_styled("10. Testing & Quality Assurance", level=1)
    add_para(
        "NexaSphere is validated by a rigorous automated test suite with 60+ test cases:"
    )
    tests = [
        "API Health & Integration Tests: Validate all 7 endpoints, KPI calculations, and question routing.",
        "Mathematical Accuracy Benchmarks: Exact-value assertions verifying profit = revenue - cost, margin formula correctness, and regional determinism (Lagos = highest revenue, Abuja = highest margin).",
        "Intent Routing Tests: Every registered label resolves to its own intent; unknown questions return helpful suggestions.",
        "Data Source Loading Tests: CSV, TSV, Excel, Parquet, and JSON round-trip validation with schema coercion and environment variable routing.",
        "DuckDB Cross-Validation: DuckDB engine results verified against Pandas engine for accuracy.",
        "Risk Scoring Monotonicity: Severity scores increase monotonically with worsening conditions.",
        "What-If Scenario Validation: Four parameterized levers produce correct naira-impact calculations.",
        "SSE Stream Contract: Server-Sent Events endpoint delivers staged progress and final structured response.",
    ]
    for t in tests:
        add_bullet(t)

    # ---- 11. DEPLOYMENT ----
    add_heading_styled("11. Deployment & Accessibility", level=1)
    add_para("NexaSphere is deployed and accessible via the following links:")

    deploy_table = doc.add_table(rows=4, cols=2)
    deploy_table.style = "Light Shading Accent 1"
    deploy_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cells = deploy_table.rows[0].cells
    cells[0].text = "Component"
    cells[1].text = "Access"
    for cell in cells:
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True

    data = [
        ("Live Web Dashboard", DEPLOYED_WEB),
        ("Mobile App (Android)", DEPLOYED_MOBILE),
        ("Source Repository", DEPLOYED_REPO),
    ]
    for i, (comp, url) in enumerate(data, 1):
        cells = deploy_table.rows[i].cells
        cells[0].text = comp
        cells[1].text = url

    add_para("")
    add_para(
        "The CI/CD pipeline (GitHub Actions) automatically tests the backend, builds and deploys "
        "the frontend to GitHub Pages, and releases the mobile web bundle on version tags.",
    )

    # ---- 12. FUTURE ROADMAP ----
    add_heading_styled("12. Future Roadmap", level=1)
    roadmap = [
        "Real-Time Data Connectors: Direct streaming from PostgreSQL, Snowflake, and BigQuery for live operational data.",
        "Automated Executive Digests: Scheduled Slack and email alerts with daily/weekly KPI summaries and anomaly highlights.",
        "Predictive ML Forecasting: Integration of time-series forecasting models for revenue prediction, demand planning, and churn analysis.",
        "Voice-Command Mobile Querying: Hands-free natural language queries via voice on the mobile app.",
        "Multi-Tenant Enterprise Deployment: Organisation-level data isolation, role-based access control, and audit logging.",
        "Custom Dashboard Builder: Drag-and-drop widget configuration for personalised executive dashboards.",
    ]
    for r in roadmap:
        add_bullet(r)

    # ---- 13. CONCLUSION ----
    add_heading_styled("13. Conclusion", level=1)
    add_para(
        "NexaSphere represents a paradigm shift in how businesses interact with their data. By "
        "enforcing a strict separation between mathematical computation and natural language "
        "generation, the system delivers something no generic AI chatbot can: 100% trustworthy "
        "numbers with executive-grade narrative intelligence."
    )
    add_para(
        "The solution is production-ready, fully tested, deployed across web and mobile, and "
        "designed for immediate business impact. NexaSphere doesn't just answer questions — it "
        "proactively identifies anomalies, models scenarios, and guides leaders toward better "
        "decisions with verified, auditable insights."
    )
    add_para(
        "Decision intelligence, perfected.",
        bold=True, size=13, color=RGBColor(79, 70, 229),
        align=WD_ALIGN_PARAGRAPH.CENTER
    )

    doc.save(str(OUT / "NexaSphere_Comprehensive_Summary.docx"))
    print("[OK] DOCX Summary generated")


def generate_summary_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak,
        Table, TableStyle, ListFlowable, ListItem
    )

    pdf_path = OUT / "NexaSphere_Comprehensive_Summary.pdf"
    doc = SimpleDocTemplate(
        str(pdf_path), pagesize=A4,
        rightMargin=2*cm, leftMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm,
    )
    styles = getSampleStyleSheet()

    title_s = ParagraphStyle('MainTitle', parent=styles['Title'], fontSize=28,
        leading=34, textColor=colors.HexColor(BRAND_INDIGO), spaceAfter=6)
    subtitle_s = ParagraphStyle('SubTitle', parent=styles['Normal'], fontSize=14,
        leading=18, textColor=colors.HexColor(BRAND_DARK), alignment=1, spaceAfter=4)
    tagline_s = ParagraphStyle('Tagline', parent=styles['Normal'], fontSize=11,
        leading=14, textColor=colors.HexColor(BRAND_INDIGO), alignment=1,
        fontName='Helvetica-Oblique', spaceAfter=12)
    h1 = ParagraphStyle('H1', parent=styles['Heading1'], fontSize=16, leading=20,
        textColor=colors.HexColor(BRAND_DARK), spaceBefore=16, spaceAfter=6,
        fontName='Helvetica-Bold')
    h2 = ParagraphStyle('H2', parent=styles['Heading2'], fontSize=13, leading=16,
        textColor=colors.HexColor(BRAND_INDIGO), spaceBefore=10, spaceAfter=4,
        fontName='Helvetica-Bold')
    body = ParagraphStyle('Body', parent=styles['Normal'], fontSize=9.5, leading=12.5,
        textColor=colors.HexColor('#334155'), spaceAfter=4)
    body_bold = ParagraphStyle('BodyBold', parent=body, fontName='Helvetica-Bold',
        textColor=colors.HexColor(BRAND_DARK))
    bullet_s = ParagraphStyle('Bullet', parent=body, leftIndent=16, bulletIndent=6,
        spaceBefore=1, spaceAfter=1)
    link_s = ParagraphStyle('Link', parent=body, textColor=colors.HexColor(BRAND_SKY),
        alignment=1, spaceAfter=8)

    story = []

    # Title page
    story.append(Spacer(1, 6*cm))
    story.append(Paragraph("NexaSphere", title_s))
    story.append(Paragraph("AI Business Intelligence Assistant", subtitle_s))
    story.append(Paragraph("Zero-Hallucination Decision Intelligence for Modern Enterprises", tagline_s))
    story.append(Spacer(1, 1*cm))
    story.append(Paragraph(
        f'<b>Live Dashboard:</b> <a href="{DEPLOYED_WEB}" color="#0EA5E9">{DEPLOYED_WEB}</a><br/>'
        f'<b>Mobile App:</b> <a href="{DEPLOYED_MOBILE}" color="#0EA5E9">{DEPLOYED_MOBILE}</a><br/>'
        f'<b>Source Code:</b> <a href="{DEPLOYED_REPO}" color="#0EA5E9">{DEPLOYED_REPO}</a>',
        link_s))
    story.append(Paragraph("Case Study 4 Solution  |  Competition Submission 2026", 
        ParagraphStyle('Date', parent=body, alignment=1, textColor=colors.HexColor('#64748B'), fontSize=9)))
    story.append(PageBreak())

    # 1. Executive Summary
    story.append(Paragraph("1. Executive Summary", h1))
    story.append(Paragraph(
        "NexaSphere is a full-stack AI Business Intelligence Assistant that transforms how "
        "executives and managers interact with operational data. By combining deterministic "
        "mathematical computation with natural language understanding, NexaSphere enables "
        "leaders to ask business questions in plain English and receive instant, accurate, "
        "chart-backed answers — with a <b>zero-hallucination guarantee</b> on every number.", body))
    story.append(Paragraph(
        "Unlike generic AI chatbots that fabricate financial figures, NexaSphere enforces a strict "
        "architectural invariant: <b>all mathematics is executed by Python's Pandas engine</b>, while the "
        "LLM (Google Gemini) serves purely as an executive interpreter that generates narrative, "
        "selects visualisations, and recommends next actions.", body))
    story.append(Paragraph(
        "The solution ships as three interconnected applications — a FastAPI backend, a React web "
        "dashboard, and an Expo mobile app — all powered by a single analytical engine answering "
        "nine critical management questions.", body))

    # 2. The Problem
    story.append(Paragraph("2. The Problem We Solve", h1))
    story.append(Paragraph(
        "Modern businesses generate vast amounts of operational data, yet extracting actionable "
        "insights remains painfully slow and error-prone:", body))
    problems = [
        "Manual BI reporting takes hours to days, delaying critical decisions.",
        "Fragmented data sources require multiple tools and expertise to cross-analyse.",
        "Generic AI chatbots confidently hallucinate financial numbers.",
        "Executives lack a single, trusted interface combining NL access with mathematical precision.",
        "Mobile access to BI insights is an afterthought.",
    ]
    for p in problems:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {p}", bullet_s))

    # 3. Our Solution
    story.append(Paragraph("3. Our Solution: NexaSphere", h1))
    story.append(Paragraph(
        "NexaSphere combines three critical capabilities:", body))
    solutions = [
        "<b>Deterministic Mathematical Engine:</b> Every number computed by Pandas with exact precision. The LLM never produces any numeric value.",
        "<b>Natural Language Interface:</b> Plain English queries answered instantly with dynamic charts, verified facts, and actionable recommendations.",
        "<b>Cross-Platform Access:</b> React web dashboard and Expo mobile app consuming the same FastAPI backend.",
        "<b>Proactive Anomaly Detection:</b> Automatically identifies and surfaces business anomalies before the executive asks.",
        "<b>What-If Simulation:</b> Model scenarios with quantified naira impact on the bottom line.",
    ]
    for s in solutions:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {s}", bullet_s))

    # 4. Target Audience
    story.append(Paragraph("4. Target Audience & Personas", h1))
    personas = [
        "<b>C-Suite Executives:</b> High-level KPI dashboards, anomaly alerts, and strategic recommendations via web and mobile.",
        "<b>Regional Managers:</b> Compare performance across regions, identify underperforming stores.",
        "<b>Marketing Managers:</b> Evaluate campaign ROI and optimise spend allocation.",
        "<b>Supply Chain Directors:</b> Monitor delivery performance, inventory stockouts, and return rates.",
    ]
    for p in personas:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {p}", bullet_s))

    # 5. How AI Is Used
    story.append(Paragraph("5. How AI Is Used — The Hybrid Architecture", h1))
    story.append(Paragraph(
        "NexaSphere's most significant innovation is its non-obvious use of AI:", body))
    ai_layers = [
        "<b>Layer 1 — Deterministic Intent Routing:</b> Keyword scoring maps questions to analysis functions. No LLM involved.",
        "<b>Layer 2 — Pandas Computation:</b> All aggregations, KPIs, chart data, and risk scores computed here.",
        "<b>Layer 3 — DuckDB Engine:</b> Columnar analytics with LRU caching, cross-validated against Pandas.",
        "<b>Layer 4 — LLM Narrative:</b> Gemini generates summaries, interpretations, and recommendations — only after numbers exist.",
        "<b>Layer 5 — Structured Response:</b> Merges all components into StructuredBIResponse JSON.",
    ]
    for l in ai_layers:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {l}", bullet_s))
    story.append(Paragraph(
        "<b>Guarantee:</b> Zero mathematical hallucination, complete auditability, deterministic "
        "reproducibility, and graceful degradation with offline fallback.", body_bold))

    # 6. Core Use Cases
    story.append(Paragraph("6. Core Use Cases & Supported Insights", h1))
    queries = [
        "Revenue & Profit Drivers", "Growth vs Profitability", "Return Anomalies",
        "Campaign ROI", "Inventory Health", "Delivery Performance",
        "Customer Segments", "Employee Performance", "Target Attainment",
    ]
    for q in queries:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {q}", bullet_s))

    # 7. Technical Architecture
    story.append(Paragraph("7. Technical Architecture", h1))
    layers = [
        "<b>Presentation:</b> React 18 (Vite, Recharts) on :3030; Expo React Native on :8081.",
        "<b>API:</b> FastAPI on :5050 with 7 endpoints including SSE streaming.",
        "<b>Orchestration:</b> Gemini LLM wrapper with offline fallback.",
        "<b>Analysis:</b> Nine Pandas functions + DuckDB engine + risk scoring + what-if simulation.",
        "<b>Data:</b> Source resolution for CSV/TSV/Excel/Parquet/JSON with schema validation.",
    ]
    for l in layers:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {l}", bullet_s))

    # 8. Innovation Layer
    story.append(Paragraph("8. Innovation Layer", h1))
    innovations = [
        "DuckDB Engine: Sub-5ms columnar queries with LRU caching.",
        "Risk Scoring: Deterministic 0-100 severity with urgency bands.",
        "What-If Simulation: Four parametrized scenario levers with naira impact.",
        "SSE Streaming: Staged progress events for real-time feedback.",
        "Proactive Anomaly Insights: Live Pandas-driven anomaly detection.",
    ]
    for inv in innovations:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {inv}", bullet_s))

    # 9. Business Impact
    story.append(Paragraph("9. Business Impact & Value Proposition", h1))
    impacts = [
        "<b>100% Mathematical Accuracy:</b> Zero hallucination risk.",
        "<b>~80% Reporting Latency Reduction:</b> Hours → seconds.",
        "<b>Proactive Decision-Making:</b> Anomaly alerts before board meetings.",
        "<b>Unified Analytics:</b> One assistant, nine questions, web and mobile.",
        "<b>Graceful Degradation:</b> Works offline with fallback narratives.",
    ]
    for imp in impacts:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {imp}", bullet_s))

    # 10. Testing
    story.append(Paragraph("10. Testing & Quality Assurance", h1))
    tests = [
        "60+ automated test cases across 5 test files.",
        "Mathematical accuracy benchmarks with exact-value assertions.",
        "Intent routing validation for all 9 registered analyses.",
        "Data source loading tests (CSV, TSV, Excel, Parquet, JSON).",
        "DuckDB cross-validation against Pandas engine.",
        "Risk scoring monotonicity and what-if scenario verification.",
    ]
    for t in tests:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {t}", bullet_s))

    # 11. Deployment
    story.append(Paragraph("11. Deployment & Accessibility", h1))
    deploy_data = [
        ['Component', 'URL'],
        ['Live Web Dashboard', DEPLOYED_WEB],
        ['Mobile App', DEPLOYED_MOBILE],
        ['Source Repository', DEPLOYED_REPO],
    ]
    t = Table(deploy_data, colWidths=[120, 350])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(BRAND_INDIGO)),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 8.5),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#F8FAFC'), colors.white]),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(t)

    # 12. Future Roadmap
    story.append(Paragraph("12. Future Roadmap", h1))
    roadmap = [
        "Real-Time Data Connectors: PostgreSQL, Snowflake, BigQuery streaming.",
        "Automated Executive Digests: Scheduled Slack and email alerts.",
        "Predictive ML Forecasting: Revenue prediction and demand planning.",
        "Voice-Command Mobile Querying: Hands-free natural language access.",
        "Multi-Tenant Enterprise Deployment: RBAC and audit logging.",
        "Custom Dashboard Builder: Drag-and-drop widget configuration.",
    ]
    for r in roadmap:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {r}", bullet_s))

    # 13. Conclusion
    story.append(Paragraph("13. Conclusion", h1))
    story.append(Paragraph(
        "NexaSphere represents a paradigm shift in how businesses interact with their data. "
        "By enforcing a strict separation between mathematical computation and natural language "
        "generation, the system delivers something no generic AI chatbot can: <b>100% trustworthy "
        "numbers with executive-grade narrative intelligence</b>.", body))
    story.append(Paragraph(
        "Decision intelligence, perfected.",
        ParagraphStyle('Conclusion', parent=body, fontSize=13, leading=16,
            textColor=colors.HexColor(BRAND_INDIGO), alignment=1,
            fontName='Helvetica-Bold', spaceBefore=12)))

    doc.build(story)
    print("[OK] PDF Summary generated")


# ---------------------------------------------------------------------------
# 2. COMPREHENSIVE SLIDES PRESENTATION (.pptx and .pdf)
# ---------------------------------------------------------------------------
def generate_slides_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK = RGBColor(15, 23, 42)
    INDIGO = RGBColor(79, 70, 229)
    SKY = RGBColor(14, 165, 233)
    EMERALD = RGBColor(16, 185, 129)
    AMBER = RGBColor(245, 158, 11)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)
    SLATE_600 = RGBColor(71, 85, 105)
    SLATE_300 = RGBColor(203, 213, 225)

    def add_bg(slide, color=DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape_bg(slide, color=DARK):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return tf

    def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                        color=WHITE, bullet_char="\u2022"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"{bullet_char}  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Calibri"
            p.space_after = Pt(6)
        return tf

    def add_accent_bar(slide, left, top, width, height, color=INDIGO):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # ========== SLIDE 1: TITLE ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_shape_bg(slide1, DARK)
    add_accent_bar(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.08), INDIGO)

    add_text_box(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "NexaSphere", font_size=54, color=INDIGO, bold=True)
    add_text_box(slide1, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
                 "AI Business Intelligence Assistant", font_size=32, color=WHITE, bold=True)
    add_text_box(slide1, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
                 "Zero-Hallucination Decision Intelligence for Modern Enterprises",
                 font_size=18, color=SKY, bold=False)
    add_accent_bar(slide1, Inches(1), Inches(4.4), Inches(3), Inches(0.04), SKY)

    links_text = (f"Live Dashboard: {DEPLOYED_WEB}\n"
                  f"Mobile App: {DEPLOYED_MOBILE}\n"
                  f"Source Code: {DEPLOYED_REPO}")
    add_text_box(slide1, Inches(1), Inches(5.0), Inches(11), Inches(1.2),
                 links_text, font_size=13, color=SLATE_300)

    add_text_box(slide1, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "Case Study 4 Solution  |  Competition Submission 2026",
                 font_size=12, color=SLATE_300, alignment=PP_ALIGN.LEFT)

    # ========== SLIDE 2: THE PROBLEM ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide2, DARK)
    add_accent_bar(slide2, Inches(0), Inches(0), Inches(13.333), Inches(0.08), AMBER)

    add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "The Problem We Solve", font_size=36, color=AMBER, bold=True)
    add_accent_bar(slide2, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), AMBER)

    problems = [
        "Manual BI reporting takes hours to days — decisions are delayed.",
        "Fragmented data across sales, returns, inventory, campaigns, and deliveries.",
        "Generic AI chatbots hallucinate financial numbers — eroding trust.",
        "No single trusted interface combining NL access with mathematical precision.",
        "Mobile BI access is an afterthought — leaders disconnected on the go.",
    ]
    add_bullet_list(slide2, Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                    problems, font_size=20, color=WHITE)

    add_text_box(slide2, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
                 "The result: slower decisions, missed revenue opportunities, and erosion of trust in AI.",
                 font_size=16, color=AMBER, bold=True)

    # ========== SLIDE 3: THE SOLUTION ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide3, DARK)
    add_accent_bar(slide3, Inches(0), Inches(0), Inches(13.333), Inches(0.08), EMERALD)

    add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Our Solution: NexaSphere", font_size=36, color=EMERALD, bold=True)
    add_accent_bar(slide3, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), EMERALD)

    solutions = [
        "Deterministic Mathematical Engine — Pandas computes every number. Zero exceptions.",
        "Natural Language Interface — Ask questions in plain English, get instant chart-backed answers.",
        "Cross-Platform Access — React web dashboard + Expo mobile app, one backend.",
        "Proactive Anomaly Detection — Issues surfaced before you even ask.",
        "What-If Simulation — Model scenarios with quantified naira impact.",
    ]
    add_bullet_list(slide3, Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                    solutions, font_size=20, color=WHITE)

    add_text_box(slide3, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
                 "100% mathematical accuracy  |  Executive-grade narrative  |  Web + Mobile",
                 font_size=16, color=EMERALD, bold=True)

    # ========== SLIDE 4: HYBRID ARCHITECTURE ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide4, DARK)
    add_accent_bar(slide4, Inches(0), Inches(0), Inches(13.333), Inches(0.08), INDIGO)

    add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "The Hybrid Architecture — How AI Is Used", font_size=36, color=INDIGO, bold=True)
    add_accent_bar(slide4, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), INDIGO)

    layers = [
        "1. Deterministic Intent Routing — Keyword scoring maps questions to analysis functions.",
        "2. Pandas Computation Engine — All KPIs, aggregations, chart data computed here.",
        "3. DuckDB Analytical Engine — Columnar analytics, sub-5ms queries, LRU caching.",
        "4. LLM Narrative Layer — Gemini generates narrative ONLY after numbers exist.",
        "5. Structured Response — Merges everything into StructuredBIResponse JSON.",
    ]
    add_bullet_list(slide4, Inches(1), Inches(1.6), Inches(11), Inches(4),
                    layers, font_size=18, color=WHITE)

    add_text_box(slide4, Inches(1), Inches(5.5), Inches(11), Inches(1),
                 "Key Invariant: The AI NEVER executes math. Every number originates from Pandas. "
                 "The LLM only writes prose.",
                 font_size=16, color=INDIGO, bold=True)

    # ========== SLIDE 5: TARGET AUDIENCE ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide5, DARK)
    add_accent_bar(slide5, Inches(0), Inches(0), Inches(13.333), Inches(0.08), SKY)

    add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Who It Serves", font_size=36, color=SKY, bold=True)
    add_accent_bar(slide5, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04), SKY)

    personas = [
        "C-Suite Executives — KPI dashboards, anomaly alerts, strategic recommendations.",
        "Regional Managers — Cross-region performance comparison, store drill-downs.",
        "Marketing Managers — Campaign ROI evaluation, spend optimisation.",
        "Supply Chain Directors — Delivery monitoring, inventory health, return analysis.",
    ]
    add_bullet_list(slide5, Inches(1), Inches(1.6), Inches(11), Inches(3.5),
                    personas, font_size=19, color=WHITE)

    add_text_box(slide5, Inches(0.8), Inches(5.2), Inches(11), Inches(0.8),
                 "The Nine Core Management Questions", font_size=22, color=SKY, bold=True)
    queries_text = ("Revenue & Profit Drivers  |  Growth vs Profitability  |  Return Anomalies  |  "
                    "Campaign ROI  |  Inventory Health  |  Delivery Performance  |  "
                    "Customer Segments  |  Employee Performance  |  Target Attainment")
    add_text_box(slide5, Inches(1), Inches(5.9), Inches(11), Inches(1),
                 queries_text, font_size=14, color=SLATE_300)

    # ========== SLIDE 6: INNOVATION LAYER ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide6, DARK)
    add_accent_bar(slide6, Inches(0), Inches(0), Inches(13.333), Inches(0.08), AMBER)

    add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Innovation Layer", font_size=36, color=AMBER, bold=True)
    add_accent_bar(slide6, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04), AMBER)

    innovations = [
        "DuckDB Engine — In-process columnar analytics, sub-5ms queries, LRU caching.",
        "Risk Scoring — Deterministic 0-100 severity with urgency bands per analysis.",
        "What-If Simulation — Four scenario levers with quantified naira impact.",
        "SSE Streaming — Staged progress events for real-time user feedback.",
        "Proactive Anomaly Insights — Live Pandas-driven anomaly detection.",
        "Structured Response Contract — Facts, interpretations, recommendations segregated.",
    ]
    add_bullet_list(slide6, Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                    innovations, font_size=19, color=WHITE)

    # ========== SLIDE 7: BUSINESS IMPACT ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide7, DARK)
    add_accent_bar(slide7, Inches(0), Inches(0), Inches(13.333), Inches(0.08), EMERALD)

    add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Business Impact & Value Proposition", font_size=36, color=EMERALD, bold=True)
    add_accent_bar(slide7, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), EMERALD)

    impacts = [
        "100% Mathematical Accuracy — Zero hallucination risk on every figure.",
        "~80% Reporting Latency Reduction — Hours of manual work → seconds.",
        "Proactive Decision-Making — Anomaly alerts before board meetings.",
        "Unified Analytics Platform — One assistant, nine questions, web + mobile.",
        "Graceful Degradation — Works offline with deterministic fallback narratives.",
        "Zero Infrastructure Overhead — No Docker/K8s required for MVP.",
    ]
    add_bullet_list(slide7, Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                    impacts, font_size=19, color=WHITE)

    # ========== SLIDE 8: TESTING ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide8, DARK)
    add_accent_bar(slide8, Inches(0), Inches(0), Inches(13.333), Inches(0.08), INDIGO)

    add_text_box(slide8, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Testing & Quality Assurance", font_size=36, color=INDIGO, bold=True)
    add_accent_bar(slide8, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), INDIGO)

    tests = [
        "60+ automated test cases across 5 comprehensive test files.",
        "Mathematical accuracy benchmarks with exact-value assertions.",
        "Intent routing validation for all 9 registered analyses.",
        "Data source loading tests (CSV, TSV, Excel, Parquet, JSON).",
        "DuckDB cross-validation against Pandas engine.",
        "Risk scoring monotonicity and what-if scenario verification.",
    ]
    add_bullet_list(slide8, Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                    tests, font_size=19, color=WHITE)

    # ========== SLIDE 9: DEPLOYMENT ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide9, DARK)
    add_accent_bar(slide9, Inches(0), Inches(0), Inches(13.333), Inches(0.08), SKY)

    add_text_box(slide9, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Deployment & Accessibility", font_size=36, color=SKY, bold=True)
    add_accent_bar(slide9, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), SKY)

    deploy_items = [
        f"Live Web Dashboard: {DEPLOYED_WEB}",
        f"Mobile App (Android): {DEPLOYED_MOBILE}",
        f"Source Repository: {DEPLOYED_REPO}",
    ]
    add_bullet_list(slide9, Inches(1), Inches(1.8), Inches(11), Inches(2.5),
                    deploy_items, font_size=18, color=WHITE)

    add_text_box(slide9, Inches(0.8), Inches(4.0), Inches(11), Inches(0.8),
                 "CI/CD Pipeline (GitHub Actions)", font_size=22, color=SKY, bold=True)
    cicd = [
        "Automated backend testing on every push/PR.",
        "Frontend auto-deployed to GitHub Pages on main branch.",
        "Mobile web bundle released as GitHub Release on version tags.",
        "Zero repository secrets required — GEMINI_API_KEY optional at runtime.",
    ]
    add_bullet_list(slide9, Inches(1), Inches(4.7), Inches(11), Inches(2.5),
                    cicd, font_size=17, color=WHITE)

    # ========== SLIDE 10: TECH STACK ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide10, DARK)
    add_accent_bar(slide10, Inches(0), Inches(0), Inches(13.333), Inches(0.08), INDIGO)

    add_text_box(slide10, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Technology Stack", font_size=36, color=INDIGO, bold=True)
    add_accent_bar(slide10, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04), INDIGO)

    stack = [
        "Backend: Python 3.10+, FastAPI, Pandas, DuckDB, Google Gemini API.",
        "Web: React 18, Vite, Recharts, lucide-react, CSS Modules.",
        "Mobile: Expo SDK 50, React Native 0.73, gifted-charts.",
        "Testing: pytest with accuracy benchmarks and data validation.",
        "CI/CD: GitHub Actions (test → deploy → release).",
        "Data: CSV, TSV, Excel, Parquet, JSON with schema validation.",
    ]
    add_bullet_list(slide10, Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                    stack, font_size=19, color=WHITE)

    # ========== SLIDE 11: FUTURE ROADMAP ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide11, DARK)
    add_accent_bar(slide11, Inches(0), Inches(0), Inches(13.333), Inches(0.08), AMBER)

    add_text_box(slide11, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Future Roadmap", font_size=36, color=AMBER, bold=True)
    add_accent_bar(slide11, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04), AMBER)

    roadmap = [
        "Real-Time Data Connectors — PostgreSQL, Snowflake, BigQuery streaming.",
        "Automated Executive Digests — Scheduled Slack and email alerts.",
        "Predictive ML Forecasting — Revenue prediction and demand planning.",
        "Voice-Command Mobile Querying — Hands-free natural language access.",
        "Multi-Tenant Enterprise Deployment — RBAC and audit logging.",
        "Custom Dashboard Builder — Drag-and-drop widget configuration.",
    ]
    add_bullet_list(slide11, Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                    roadmap, font_size=19, color=WHITE)

    # ========== SLIDE 12: CLOSING ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide12, DARK)
    add_accent_bar(slide12, Inches(0), Inches(0), Inches(13.333), Inches(0.08), INDIGO)

    add_text_box(slide12, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "Decision Intelligence, Perfected.", font_size=44, color=INDIGO,
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide12, Inches(1), Inches(2.8), Inches(11), Inches(1.5),
                 "NexaSphere delivers what no generic AI chatbot can:\n"
                 "100% trustworthy numbers with executive-grade narrative intelligence.",
                 font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide12, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.04), SKY)

    closing_links = (f"Live Dashboard: {DEPLOYED_WEB}\n"
                     f"Mobile App: {DEPLOYED_MOBILE}\n"
                     f"Source Code: {DEPLOYED_REPO}")
    add_text_box(slide12, Inches(1), Inches(4.8), Inches(11), Inches(1.5),
                 closing_links, font_size=14, color=SLATE_300, alignment=PP_ALIGN.CENTER)

    add_text_box(slide12, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "Thank you  |  NexaSphere AI BI Assistant",
                 font_size=13, color=SLATE_300, alignment=PP_ALIGN.CENTER)

    prs.save(str(OUT / "NexaSphere_Presentation_Slides.pptx"))
    print("[OK] PPTX Slides generated")


def generate_slides_pdf():
    """Convert the PPTX to PDF using reportlab (re-create as PDF pages)."""
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak

    pdf_path = OUT / "NexaSphere_Presentation_Slides.pdf"
    W, H = landscape(A4)
    doc = SimpleDocTemplate(str(pdf_path), pagesize=landscape(A4),
        rightMargin=1.5*cm, leftMargin=1.5*cm, topMargin=1*cm, bottomMargin=1*cm)
    styles = getSampleStyleSheet()

    dark = colors.HexColor(BRAND_DARK)
    indigo = colors.HexColor(BRAND_INDIGO)
    sky = colors.HexColor(BRAND_SKY)
    emerald = colors.HexColor(BRAND_EMERALD)
    amber = colors.HexColor(BRAND_AMBER)
    white = colors.white
    slate = colors.HexColor('#94A3B8')

    title_s = ParagraphStyle('ST', fontSize=36, leading=44, textColor=indigo,
        fontName='Helvetica-Bold', alignment=1, spaceAfter=12)
    subtitle_s = ParagraphStyle('SS', fontSize=20, leading=26, textColor=white,
        fontName='Helvetica-Bold', alignment=1, spaceAfter=8)
    tagline_s = ParagraphStyle('TG', fontSize=14, leading=18, textColor=sky,
        fontName='Helvetica-Oblique', alignment=1, spaceAfter=20)
    h1 = ParagraphStyle('H1', fontSize=28, leading=34, textColor=indigo,
        fontName='Helvetica-Bold', spaceBefore=8, spaceAfter=12)
    h1_amber = ParagraphStyle('H1A', parent=h1, textColor=amber)
    h1_emerald = ParagraphStyle('H1E', parent=h1, textColor=emerald)
    h1_sky = ParagraphStyle('H1S', parent=h1, textColor=sky)
    body = ParagraphStyle('B', fontSize=13, leading=17, textColor=white,
        fontName='Helvetica', spaceAfter=4)
    bullet_s = ParagraphStyle('BL', parent=body, leftIndent=20, bulletIndent=8,
        spaceBefore=2, spaceAfter=3)
    link_s = ParagraphStyle('LK', fontSize=11, leading=14, textColor=slate,
        alignment=1, spaceAfter=6)
    accent_s = ParagraphStyle('AC', fontSize=14, leading=18, textColor=indigo,
        fontName='Helvetica-Bold', alignment=1, spaceBefore=12)

    story = []

    # Slide 1: Title
    story.append(Spacer(1, 4*cm))
    story.append(Paragraph("NexaSphere", title_s))
    story.append(Paragraph("AI Business Intelligence Assistant", subtitle_s))
    story.append(Paragraph("Zero-Hallucination Decision Intelligence for Modern Enterprises", tagline_s))
    story.append(Paragraph(
        f"Live Dashboard: {DEPLOYED_WEB}<br/>"
        f"Mobile App: {DEPLOYED_MOBILE}<br/>"
        f"Source Code: {DEPLOYED_REPO}", link_s))
    story.append(Paragraph("Case Study 4 Solution  |  Competition Submission 2026",
        ParagraphStyle('D', parent=link_s, textColor=slate)))
    story.append(PageBreak())

    # Slide 2: Problem
    story.append(Paragraph("The Problem We Solve", h1_amber))
    probs = [
        "Manual BI reporting takes hours to days — decisions are delayed.",
        "Fragmented data across sales, returns, inventory, campaigns, and deliveries.",
        "Generic AI chatbots hallucinate financial numbers — eroding trust.",
        "No single trusted interface combining NL access with mathematical precision.",
        "Mobile BI access is an afterthought — leaders disconnected on the go.",
    ]
    for p in probs:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {p}", bullet_s))
    story.append(Paragraph("The result: slower decisions, missed revenue opportunities, and erosion of trust in AI.",
        accent_s))
    story.append(PageBreak())

    # Slide 3: Solution
    story.append(Paragraph("Our Solution: NexaSphere", h1_emerald))
    sols = [
        "<b>Deterministic Mathematical Engine</b> — Pandas computes every number. Zero exceptions.",
        "<b>Natural Language Interface</b> — Ask questions in plain English, get instant chart-backed answers.",
        "<b>Cross-Platform Access</b> — React web dashboard + Expo mobile app, one backend.",
        "<b>Proactive Anomaly Detection</b> — Issues surfaced before you even ask.",
        "<b>What-If Simulation</b> — Model scenarios with quantified naira impact.",
    ]
    for s in sols:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {s}", bullet_s))
    story.append(Paragraph("100% mathematical accuracy  |  Executive-grade narrative  |  Web + Mobile",
        ParagraphStyle('IMP', parent=accent_s, textColor=emerald)))
    story.append(PageBreak())

    # Slide 4: Architecture
    story.append(Paragraph("The Hybrid Architecture — How AI Is Used", h1))
    layers = [
        "<b>1. Deterministic Intent Routing</b> — Keyword scoring maps questions to analysis functions.",
        "<b>2. Pandas Computation Engine</b> — All KPIs, aggregations, chart data computed here.",
        "<b>3. DuckDB Analytical Engine</b> — Columnar analytics, sub-5ms queries, LRU caching.",
        "<b>4. LLM Narrative Layer</b> — Gemini generates narrative ONLY after numbers exist.",
        "<b>5. Structured Response</b> — Merges everything into StructuredBIResponse JSON.",
    ]
    for l in layers:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {l}", bullet_s))
    story.append(Paragraph("Key Invariant: The AI NEVER executes math. Every number originates from Pandas.",
        accent_s))
    story.append(PageBreak())

    # Slide 5: Audience
    story.append(Paragraph("Who It Serves", h1_sky))
    personas = [
        "<b>C-Suite Executives</b> — KPI dashboards, anomaly alerts, strategic recommendations.",
        "<b>Regional Managers</b> — Cross-region performance comparison, store drill-downs.",
        "<b>Marketing Managers</b> — Campaign ROI evaluation, spend optimisation.",
        "<b>Supply Chain Directors</b> — Delivery monitoring, inventory health, return analysis.",
    ]
    for p in personas:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {p}", bullet_s))
    story.append(Spacer(1, 8))
    story.append(Paragraph("The Nine Core Management Questions", ParagraphStyle('Q', parent=h1, textColor=sky, fontSize=20)))
    story.append(Paragraph(
        "Revenue & Profit Drivers  |  Growth vs Profitability  |  Return Anomalies  |  "
        "Campaign ROI  |  Inventory Health  |  Delivery Performance  |  "
        "Customer Segments  |  Employee Performance  |  Target Attainment",
        ParagraphStyle('QT', parent=body, fontSize=11, textColor=slate)))
    story.append(PageBreak())

    # Slide 6: Innovation
    story.append(Paragraph("Innovation Layer", h1_amber))
    inno = [
        "<b>DuckDB Engine</b> — In-process columnar analytics, sub-5ms queries, LRU caching.",
        "<b>Risk Scoring</b> — Deterministic 0-100 severity with urgency bands per analysis.",
        "<b>What-If Simulation</b> — Four scenario levers with quantified naira impact.",
        "<b>SSE Streaming</b> — Staged progress events for real-time user feedback.",
        "<b>Proactive Anomaly Insights</b> — Live Pandas-driven anomaly detection.",
        "<b>Structured Response Contract</b> — Facts, interpretations, recommendations segregated.",
    ]
    for i in inno:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {i}", bullet_s))
    story.append(PageBreak())

    # Slide 7: Impact
    story.append(Paragraph("Business Impact & Value Proposition", h1_emerald))
    imps = [
        "<b>100% Mathematical Accuracy</b> — Zero hallucination risk on every figure.",
        "<b>~80% Reporting Latency Reduction</b> — Hours of manual work to seconds.",
        "<b>Proactive Decision-Making</b> — Anomaly alerts before board meetings.",
        "<b>Unified Analytics Platform</b> — One assistant, nine questions, web + mobile.",
        "<b>Graceful Degradation</b> — Works offline with deterministic fallback narratives.",
        "<b>Zero Infrastructure Overhead</b> — No Docker/K8s required for MVP.",
    ]
    for i in imps:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {i}", bullet_s))
    story.append(PageBreak())

    # Slide 8: Testing
    story.append(Paragraph("Testing & Quality Assurance", h1))
    tests = [
        "60+ automated test cases across 5 comprehensive test files.",
        "Mathematical accuracy benchmarks with exact-value assertions.",
        "Intent routing validation for all 9 registered analyses.",
        "Data source loading tests (CSV, TSV, Excel, Parquet, JSON).",
        "DuckDB cross-validation against Pandas engine.",
        "Risk scoring monotonicity and what-if scenario verification.",
    ]
    for t in tests:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {t}", bullet_s))
    story.append(PageBreak())

    # Slide 9: Deployment
    story.append(Paragraph("Deployment & Accessibility", h1_sky))
    story.append(Paragraph(f"<b>Live Web Dashboard:</b> {DEPLOYED_WEB}", body))
    story.append(Paragraph(f"<b>Mobile App:</b> {DEPLOYED_MOBILE}", body))
    story.append(Paragraph(f"<b>Source Repository:</b> {DEPLOYED_REPO}", body))
    story.append(Spacer(1, 12))
    story.append(Paragraph("CI/CD Pipeline (GitHub Actions)", ParagraphStyle('CI', parent=h1, textColor=sky, fontSize=20)))
    cicd = [
        "Automated backend testing on every push/PR.",
        "Frontend auto-deployed to GitHub Pages on main branch.",
        "Mobile web bundle released as GitHub Release on version tags.",
        "Zero repository secrets required.",
    ]
    for c in cicd:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {c}", bullet_s))
    story.append(PageBreak())

    # Slide 10: Tech Stack
    story.append(Paragraph("Technology Stack", h1))
    stack = [
        "<b>Backend:</b> Python 3.10+, FastAPI, Pandas, DuckDB, Google Gemini API.",
        "<b>Web:</b> React 18, Vite, Recharts, lucide-react, CSS Modules.",
        "<b>Mobile:</b> Expo SDK 50, React Native 0.73, gifted-charts.",
        "<b>Testing:</b> pytest with accuracy benchmarks and data validation.",
        "<b>CI/CD:</b> GitHub Actions (test, deploy, release).",
        "<b>Data:</b> CSV, TSV, Excel, Parquet, JSON with schema validation.",
    ]
    for s in stack:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {s}", bullet_s))
    story.append(PageBreak())

    # Slide 11: Roadmap
    story.append(Paragraph("Future Roadmap", h1_amber))
    roadmap = [
        "Real-Time Data Connectors — PostgreSQL, Snowflake, BigQuery streaming.",
        "Automated Executive Digests — Scheduled Slack and email alerts.",
        "Predictive ML Forecasting — Revenue prediction and demand planning.",
        "Voice-Command Mobile Querying — Hands-free natural language access.",
        "Multi-Tenant Enterprise Deployment — RBAC and audit logging.",
        "Custom Dashboard Builder — Drag-and-drop widget configuration.",
    ]
    for r in roadmap:
        story.append(Paragraph(f"<bullet>&bull;</bullet> {r}", bullet_s))
    story.append(PageBreak())

    # Slide 12: Closing
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph("Decision Intelligence, Perfected.", title_s))
    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph(
        "NexaSphere delivers what no generic AI chatbot can:<br/>"
        "100% trustworthy numbers with executive-grade narrative intelligence.",
        subtitle_s))
    story.append(Spacer(1, 1*cm))
    story.append(Paragraph(
        f"Live Dashboard: {DEPLOYED_WEB}<br/>"
        f"Mobile App: {DEPLOYED_MOBILE}<br/>"
        f"Source Code: {DEPLOYED_REPO}", link_s))
    story.append(Paragraph("Thank you  |  NexaSphere AI BI Assistant",
        ParagraphStyle('TH', parent=link_s, textColor=slate, fontSize=12)))

    # Build with dark background
    def on_page(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(dark)
        canvas.rect(0, 0, W, H, fill=1, stroke=0)
        # Top accent bar
        canvas.setFillColor(indigo)
        canvas.rect(0, H - 6, W, 6, fill=1, stroke=0)
        canvas.restoreState()

    doc.build(story, onFirstPage=on_page, onLaterPages=on_page)
    print("[OK] PDF Slides generated")


# ---------------------------------------------------------------------------
# 3. DEMO VIDEO (.mp4) with male Nigerian voice narrative
# ---------------------------------------------------------------------------
def generate_demo_video():
    import asyncio
    import edge_tts
    from PIL import Image, ImageDraw, ImageFont
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
    import tempfile

    VOICE = "en-NG-AbeoNeural"  # Male Nigerian English voice

    # ---- 2-minute narrative script (120 seconds) ----
    segments = [
        {
            "text": (
                "Executives today face a critical dilemma. "
                "They either wait days for manual business intelligence reports, "
                "or they rely on AI chatbots that hallucinate financial numbers. "
                "Today, we solve this problem permanently. "
                "Welcome to NexaSphere, the AI Business Intelligence Assistant "
                "that combines one hundred percent exact mathematical precision "
                "with executive level reasoning."
            ),
            "shot": "01_executive_dashboard.png",
            "duration": 18,
        },
        {
            "text": (
                "NexaSphere runs on a revolutionary hybrid engine. "
                "Python's Pandas engine strictly computes every single number. "
                "The AI language model acts purely as an executive interpreter. "
                "Zero numerical hallucinations. Guaranteed. "
                "Watch as we ask: Which region generates the highest revenue and profit margin?"
            ),
            "shot": "01_executive_dashboard.png",
            "duration": 16,
        },
        {
            "text": (
                "In seconds, the system calculates exact figures "
                "and renders dynamic visualizations. "
                "It isolates verified facts from actionable recommendations, "
                "then suggests the next logical business questions "
                "through interactive follow up chips."
            ),
            "shot": "02_regional_revenue_chart.png",
            "duration": 14,
        },
        {
            "text": (
                "Next, let's examine return rate anomalies. "
                "NexaSphere identifies products with unusually high return rates, "
                "computes severity scores, and recommends corrective actions "
                "all within a single structured response."
            ),
            "shot": "03_return_rates_analysis.png",
            "duration": 14,
        },
        {
            "text": (
                "Marketing managers can instantly evaluate campaign return on investment. "
                "NexaSphere ranks campaigns by ROI, visualizes spend allocation, "
                "and quantifies the impact of budget reallocation scenarios."
            ),
            "shot": "04_marketing_roi_donut_chart.png",
            "duration": 13,
        },
        {
            "text": (
                "The performance drill down page provides deep operational visibility. "
                "Executives can explore all nine business dimensions, "
                "export data to CSV, and make data driven decisions with confidence."
            ),
            "shot": "05_performance_drilldown.png",
            "duration": 12,
        },
        {
            "text": (
                "On the go? The Expo mobile app gives executives "
                "instant access to the exact same backend engine. "
                "Tap a quick question chip, inspect insights, "
                "and receive immediate operational guidance right from your phone."
            ),
            "shot": "06_mobile_dashboard.png",
            "duration": 14,
        },
        {
            "text": (
                "The mobile assistant features intelligent question suggestions. "
                "Tap any chip to get instant answers with verified facts, "
                "interpretations, and strategic recommendations."
            ),
            "shot": "07_mobile_assistant_chips.png",
            "duration": 10,
        },
        {
            "text": (
                "Verified by rigorous automated accuracy tests, "
                "NexaSphere delivers zero hallucination guarantees, "
                "cuts reporting latency by eighty percent, "
                "and empowers leaders to act immediately. "
                "NexaSphere. Decision intelligence, perfected. "
                "Try it now at the live dashboard link."
            ),
            "shot": "08_mobile_answer.png",
            "duration": 12,
        },
    ]

    total_duration = sum(s["duration"] for s in segments)
    print(f"[video] Script total: {total_duration}s")

    # Ensure we are at ~120s, adjust last segment
    diff = 120 - total_duration
    if diff != 0:
        segments[-1]["duration"] += diff
        total_duration = 120
    print(f"[video] Adjusted to: {total_duration}s")

    tmp_dir = Path(tempfile.mkdtemp(prefix="nexasphere_video_"))

    # Step 1: Generate TTS audio for each segment
    print("[video] Generating TTS audio segments...")
    audio_files = []

    async def generate_audio():
        for i, seg in enumerate(segments):
            audio_path = str(tmp_dir / f"seg_{i:02d}.mp3")
            communicate = edge_tts.Communicate(seg["text"], VOICE, rate="-5%")
            await communicate.save(audio_path)
            audio_files.append(audio_path)
            print(f"  audio segment {i+1}/{len(segments)} done")

    asyncio.run(generate_audio())

    # Step 2: Create visual segments from screenshots
    print("[video] Creating visual segments...")
    video_clips = []

    for i, seg in enumerate(segments):
        shot_path = SHOTS / seg["shot"]
        if not shot_path.exists():
            print(f"  WARNING: {shot_path} not found, using placeholder")
            # Create a placeholder
            img = Image.new("RGB", (1920, 1080), (15, 23, 42))
            draw = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 48)
            except OSError:
                font = ImageFont.load_default()
            draw.text((960, 540), "NexaSphere", fill=(79, 70, 229), font=font, anchor="mm")
            placeholder_path = str(tmp_dir / f"placeholder_{i}.png")
            img.save(placeholder_path)
            shot_path = Path(placeholder_path)

        # Open and resize to 1920x1080
        img = Image.open(shot_path).convert("RGB")
        img = img.resize((1920, 1080), Image.LANCZOS)

        # Add text overlay at bottom
        draw = ImageDraw.Draw(img)
        # Semi-transparent bottom bar
        overlay = Image.new("RGBA", (1920, 120), (15, 23, 42, 200))
        img_rgba = img.convert("RGBA")
        img_rgba.paste(overlay, (0, 960), overlay)
        img = img_rgba.convert("RGB")

        # Add text
        try:
            font = ImageFont.truetype("arial.ttf", 28)
            small_font = ImageFont.truetype("arial.ttf", 18)
        except OSError:
            font = ImageFont.load_default()
            small_font = font

        draw = ImageDraw.Draw(img)
        # Title
        draw.text((40, 975), "NexaSphere AI BI Assistant", fill=(79, 70, 229), font=font)
        # Deployed link
        draw.text((40, 1015), f"Live: {DEPLOYED_WEB}", fill=(14, 165, 233), font=small_font)
        # Segment indicator
        draw.text((1600, 990), f"{i+1}/{len(segments)}", fill=(148, 163, 184), font=small_font)

        frame_path = str(tmp_dir / f"frame_{i:02d}.png")
        img.save(frame_path)

        # Create video clip from image + audio
        audio_clip = AudioFileClip(audio_files[i])
        img_clip = ImageClip(frame_path).with_duration(audio_clip.duration)
        img_clip = img_clip.with_audio(audio_clip)
        video_clips.append(img_clip)
        print(f"  visual segment {i+1}/{len(segments)} done ({audio_clip.duration:.1f}s)")

    # Step 3: Concatenate all segments
    print("[video] Concatenating segments...")
    final_video = concatenate_videoclips(video_clips, method="compose")

    output_path = str(OUT / "NexaSphere_Demo_Video.mp4")
    final_video.write_videofile(
        output_path,
        fps=24,
        codec="libx264",
        audio_codec="aac",
        preset="medium",
        threads=4,
        logger=None,
    )

    # Cleanup
    final_video.close()
    for clip in video_clips:
        clip.close()
    shutil.rmtree(tmp_dir, ignore_errors=True)

    size_mb = os.path.getsize(output_path) / (1024 * 1024)
    print(f"[OK] Demo video generated: {output_path} ({size_mb:.1f} MB)")


# ---------------------------------------------------------------------------
# 4. UPDATE MANIFEST
# ---------------------------------------------------------------------------
def write_manifest():
    from datetime import datetime, timezone
    rows = ["# Submission Manifest", "",
            f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}", "",
            "## Comprehensive Submission Package", ""]
    for path in sorted(OUT.rglob("*")):
        if path.is_file() and path.name != "MANIFEST.md":
            size = path.stat().st_size
            pretty = f"{size / 1e6:.1f} MB" if size > 1e6 else f"{size / 1024:.0f} KB"
            rows.append(f"- `{path.relative_to(OUT).as_posix()}` ({pretty})")
    (OUT / "MANIFEST.md").write_text("\n".join(rows) + "\n", encoding="utf-8")
    print("[OK] MANIFEST.md updated")


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------
def main():
    print("=" * 60)
    print("NexaSphere Submission Artifact Generator")
    print("=" * 60)

    print("\n[1/4] Generating Comprehensive Summary (.docx)...")
    generate_summary_docx()

    print("\n[2/4] Generating Comprehensive Summary (.pdf)...")
    generate_summary_pdf()

    print("\n[3/4] Generating Presentation Slides (.pptx)...")
    generate_slides_pptx()

    print("\n[3b/4] Generating Presentation Slides (.pdf)...")
    generate_slides_pdf()

    print("\n[4/4] Generating Demo Video (.mp4) with Nigerian male voice...")
    generate_demo_video()

    print("\n[final] Updating MANIFEST.md...")
    write_manifest()

    print("\n" + "=" * 60)
    print("ALL ARTIFACTS GENERATED SUCCESSFULLY")
    print("=" * 60)
    print(f"\nOutput directory: {OUT}")
    for f in sorted(p.name for p in OUT.iterdir() if p.is_file()):
        size = (OUT / f).stat().st_size
        pretty = f"{size / 1e6:.1f} MB" if size > 1e6 else f"{size / 1024:.0f} KB"
        print(f"  {f} ({pretty})")


if __name__ == "__main__":
    main()
